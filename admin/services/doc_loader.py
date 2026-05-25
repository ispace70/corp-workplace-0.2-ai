"""
문서 로더 — 파일 포맷별 텍스트 추출

지원 포맷: PDF, DOCX, DOC, TXT, MD, HTML, HWPX, XLS, XLSX, PPT, PPTX

반환 형식:
  list[dict]: [{"text": "...", "page": 1}, ...]
  페이지/시트별로 분리하여 반환 (청킹 시 페이지 메타 보존)
"""

import os
import re
from pathlib import Path


def load_document(file_path: str) -> list[dict]:
    """
    파일 경로를 받아 페이지/시트별 텍스트 목록 반환.
    각 항목: {"text": str, "page": int}
    """
    path = Path(file_path)
    ext = path.suffix.lower()

    loaders = {
        ".pdf":  _load_pdf,
        ".docx": _load_docx,
        ".doc":  _load_docx,
        ".txt":  _load_text,
        ".md":   _load_text,
        ".html": _load_html,
        ".htm":  _load_html,
        ".hwpx": _load_hwpx,
        ".hwp":  _load_hwp,
        ".xlsx": _load_excel,
        ".xls":  _load_excel,
        ".pptx": _load_pptx,
        ".ppt":  _load_pptx,
    }

    loader = loaders.get(ext)
    if loader is None:
        raise ValueError(f"지원하지 않는 파일 형식입니다: {ext}")

    pages = loader(str(path))
    # 빈 페이지 제거 및 텍스트 정리
    return [{"text": _clean(p["text"]), "page": p["page"]}
            for p in pages if p["text"].strip()]


def _clean(text: str) -> str:
    text = re.sub(r'\n{3,}', '\n\n', text)
    text = re.sub(r'[ \t]{2,}', ' ', text)
    return text.strip()


# ── PDF ───────────────────────────────────────────────────────────────────────

def _load_pdf(path: str) -> list[dict]:
    try:
        from pypdf import PdfReader
    except ImportError:
        raise RuntimeError("pypdf 패키지가 필요합니다: pip install pypdf")

    reader = PdfReader(path)
    pages = []
    for i, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        pages.append({"text": text, "page": i})
    return pages


# ── DOCX ──────────────────────────────────────────────────────────────────────

def _load_docx(path: str) -> list[dict]:
    try:
        from docx import Document
    except ImportError:
        raise RuntimeError("python-docx 패키지가 필요합니다: pip install python-docx")

    doc = Document(path)
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    # 테이블 텍스트도 포함
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                paragraphs.append(row_text)

    # 500자 단위로 가상 페이지 분할 (DOCX는 페이지 개념 없음)
    return _split_into_pages("\n".join(paragraphs))


# ── TXT / MD ─────────────────────────────────────────────────────────────────

def _load_text(path: str) -> list[dict]:
    with open(path, encoding="utf-8", errors="ignore") as f:
        text = f.read()
    return _split_into_pages(text)


# ── HTML ─────────────────────────────────────────────────────────────────────

def _load_html(path: str) -> list[dict]:
    try:
        from bs4 import BeautifulSoup
    except ImportError:
        raise RuntimeError("beautifulsoup4 패키지가 필요합니다: pip install beautifulsoup4 lxml")

    with open(path, encoding="utf-8", errors="ignore") as f:
        soup = BeautifulSoup(f.read(), "lxml")

    for tag in soup(["script", "style", "nav", "footer", "header"]):
        tag.decompose()

    text = soup.get_text(separator="\n")
    return _split_into_pages(text)


# ── HWPX ─────────────────────────────────────────────────────────────────────

def _load_hwpx(path: str) -> list[dict]:
    import zipfile

    try:
        from bs4 import BeautifulSoup
    except ImportError:
        raise RuntimeError("beautifulsoup4 패키지가 필요합니다: pip install beautifulsoup4 lxml")

    texts = []
    try:
        with zipfile.ZipFile(path) as z:
            # HWPX 내부 섹션 XML 파일 목록
            section_files = sorted(
                [n for n in z.namelist() if re.match(r'Contents/section\d+\.xml', n)]
            )
            if not section_files:
                # 대안 경로
                section_files = sorted(
                    [n for n in z.namelist() if 'section' in n.lower() and n.endswith('.xml')]
                )

            for sf in section_files:
                with z.open(sf) as f:
                    xml = f.read().decode("utf-8", errors="ignore")
                soup = BeautifulSoup(xml, "lxml-xml")
                # 텍스트 노드 추출
                for t_tag in soup.find_all("t"):
                    txt = t_tag.get_text()
                    if txt.strip():
                        texts.append(txt)
    except Exception as e:
        raise RuntimeError(f"HWPX 파싱 오류: {e}")

    return _split_into_pages("\n".join(texts))


# ── HWP (바이너리) ────────────────────────────────────────────────────────────

def _load_hwp(path: str) -> list[dict]:
    try:
        import olefile
    except ImportError:
        # olefile 없으면 텍스트 추출 시도 (제한적)
        return [{"text": f"[HWP 파일: {Path(path).name}] — HWP 파싱을 위해 olefile 패키지를 설치하세요.", "page": 1}]

    try:
        ole = olefile.OleFileIO(path)
        if ole.exists("BodyText/Section0"):
            data = ole.openstream("BodyText/Section0").read()
            # HWP 텍스트 레코드에서 한글 텍스트 추출 (간이 파싱)
            text = _extract_hwp_text(data)
            return _split_into_pages(text)
    except Exception:
        pass
    return [{"text": f"[HWP 파일: {Path(path).name}] — 텍스트 추출 실패", "page": 1}]


def _extract_hwp_text(data: bytes) -> str:
    """HWP BodyText 섹션 바이너리에서 텍스트 추출 (간이)."""
    import struct
    texts = []
    i = 0
    while i + 4 <= len(data):
        tag_id = struct.unpack_from('<H', data, i)[0] & 0x3FF
        size   = struct.unpack_from('<I', data, i)[0] >> 20
        i += 4
        if i + size > len(data):
            break
        if tag_id == 67:  # HWPTAG_PARA_TEXT
            chunk = data[i:i+size]
            text  = chunk.decode('utf-16-le', errors='ignore')
            texts.append(text.replace('\x00', ''))
        i += size
    return '\n'.join(texts)


# ── Excel ─────────────────────────────────────────────────────────────────────

def _load_excel(path: str) -> list[dict]:
    try:
        import openpyxl
    except ImportError:
        raise RuntimeError("openpyxl 패키지가 필요합니다: pip install openpyxl")

    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    pages = []
    for sheet_idx, ws in enumerate(wb.worksheets, start=1):
        rows_text = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            row_str = " | ".join(cells).strip(" |")
            if row_str:
                rows_text.append(row_str)
        if rows_text:
            pages.append({"text": f"[시트: {ws.title}]\n" + "\n".join(rows_text), "page": sheet_idx})
    return pages or [{"text": "", "page": 1}]


# ── PPT / PPTX ────────────────────────────────────────────────────────────────

def _load_pptx(path: str) -> list[dict]:
    try:
        from pptx import Presentation
    except ImportError:
        raise RuntimeError("python-pptx 패키지가 필요합니다: pip install python-pptx")

    prs = Presentation(path)
    pages = []
    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        if texts:
            pages.append({"text": "\n".join(texts), "page": i})
    return pages or [{"text": "", "page": 1}]


# ── 공통 유틸 ─────────────────────────────────────────────────────────────────

def _split_into_pages(text: str, page_size: int = 2000) -> list[dict]:
    """긴 텍스트를 가상 페이지로 분할."""
    if not text.strip():
        return [{"text": "", "page": 1}]
    chunks = []
    for i in range(0, max(1, len(text)), page_size):
        chunks.append({"text": text[i:i+page_size], "page": len(chunks) + 1})
    return chunks
